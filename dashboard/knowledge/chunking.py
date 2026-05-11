"""Sliding-window chunking and vision-OCR document conversion.

This module provides:

- :class:`SlidingWindowChunker` — splits text into page-sized segments,
  then groups them into overlapping sliding windows. Used by both the
  ingestion pipeline and the MarkitdownReader to produce page-ranged
  chunk dicts with provenance metadata.

- :class:`VisionSlidingWindowConverter` — a custom
  ``markitdown.DocumentConverter`` that renders PDF pages to images
  via PyMuPDF, groups them into sliding windows, and sends each
  combined image to a vision LLM for structured OCR extraction. When
  no LLM client is available, ``accepts()`` returns ``False`` and
  MarkItDown falls back to the built-in ``PdfConverter`` (pdfminer).

Design points:

- **Pluggable via MarkItDown's converter registry.** The vision
  converter is registered with ``priority=-1.0`` (higher than the
  built-in ``PdfConverter`` at ``0.0``) so it takes precedence when
  an ``llm_client`` is present. No ``if/else`` branching needed in
  the Ingestor — dispatch is automatic.

- **Shared chunking logic.** Both :meth:`Ingestor._parse_file` and
  :meth:`MarkitdownReader._docs_from_text` delegate to
  :class:`SlidingWindowChunker` so the sliding-window algorithm and
  metadata shape (``page_range``, ``window_start``, ``total_pages``,
  etc.) are consistent across both entry points.

- **Parallel window processing.** PDF page windows are processed
  concurrently via ``ThreadPoolExecutor`` bounded by
  ``min(MAX_WORKERS, len(windows))``.

- **Graceful degradation.** When PyMuPDF is not installed, or when
  the LLM client is unavailable, the converter declines the file and
  MarkItDown falls back to pdfminer-based text extraction.
"""

from __future__ import annotations

import io
import logging
import uuid
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Any
from markitdown import DocumentConverter
from dashboard.knowledge.config import (
    SLIDING_WINDOW_OVERLAP,
    SLIDING_WINDOW_SIZE,
)

logger = logging.getLogger(__name__)

MAX_WORKERS = 4


class SlidingWindowChunker:
    """Chunk text using sliding windows over page-sized segments.

    Each "page" is a paragraph-bounded segment of approximately
    ``page_chars`` characters. Groups of pages are combined into
    overlapping windows — matching the pattern used for PDF vision OCR.

    Parameters
    ----------
    window_size:
        Number of pages per sliding window (default: 3).
    overlap:
        Number of pages shared between consecutive windows (default: 1).
    page_chars:
        Target character count per "page" segment (default: 1024).
    """

    def __init__(
        self,
        window_size: int = SLIDING_WINDOW_SIZE,
        overlap: int = SLIDING_WINDOW_OVERLAP,
        page_chars: int = 1024,
    ) -> None:
        if window_size < 1:
            raise ValueError("window_size must be at least 1")
        if overlap < 0:
            raise ValueError("overlap must be non-negative")
        if overlap >= window_size:
            raise ValueError("overlap must be less than window_size")
        self.window_size = window_size
        self.overlap = overlap
        self.page_chars = page_chars

    @staticmethod
    def split_into_pages(text: str, max_page_chars: int = 1024) -> list[str]:
        """Split *text* into page-sized segments on paragraph boundaries.

        Prefers double-newline paragraph boundaries; falls back to hard
        character splitting for oversized paragraphs.  Returns a non-empty
        list of non-empty strings.
        """
        if not text:
            return []
        paras = [p.strip() for p in text.split("\n\n") if p.strip()]
        if not paras:
            paras = [text.strip()]

        pages: list[str] = []
        buf = ""
        for para in paras:
            if len(buf) + len(para) + 2 <= max_page_chars:
                buf = f"{buf}\n\n{para}".strip() if buf else para
            else:
                if buf:
                    pages.append(buf)
                if len(para) > max_page_chars:
                    for i in range(0, len(para), max_page_chars):
                        chunk = para[i : i + max_page_chars].strip()
                        if chunk:
                            pages.append(chunk)
                    buf = ""
                else:
                    buf = para
        if buf:
            pages.append(buf)
        return pages if pages else [text.strip()]

    @staticmethod
    def create_sliding_windows(
        total_pages: int,
        window_size: int = SLIDING_WINDOW_SIZE,
        overlap: int = SLIDING_WINDOW_OVERLAP,
    ) -> list[tuple[int, list[int]]]:
        """Create sliding windows over page indices.

        Returns a list of ``(window_start, [page_indices])`` tuples,
        sorted by *window_start*.

        Raises :class:`ValueError` for invalid parameters.
        """
        if window_size < 1:
            raise ValueError("window_size must be at least 1")
        if overlap < 0:
            raise ValueError("overlap must be non-negative")
        if overlap >= window_size:
            raise ValueError("overlap must be less than window_size")

        windows: list[tuple[int, list[int]]] = []
        step_size = window_size - overlap

        if total_pages <= window_size:
            windows.append((0, list(range(0, total_pages))))
        else:
            i = 0
            while i < total_pages:
                window_end = min(i + window_size, total_pages)
                window_pages = list(range(i, window_end))
                windows.append((i, window_pages))
                if window_end == total_pages:
                    break
                i += step_size
                if i + window_size > total_pages and i < total_pages:
                    final_start = max(i, total_pages - window_size)
                    existing_starts = {start for start, _ in windows}
                    if final_start not in existing_starts:
                        windows.append(
                            (final_start, list(range(final_start, total_pages)))
                        )
                    break

        return windows

    def chunk(self, text: str) -> list[dict]:
        """Full pipeline: *text* → pages → windows → chunk dicts.

        Each chunk dict has ``text`` and ``metadata`` keys.  Metadata
        includes ``page_range``, ``window_start``, ``total_pages``,
        ``window_size``, ``overlap``, ``chunk_index``, and
        ``total_chunks``.
        """
        if not text or not text.strip():
            return []

        pages = self.split_into_pages(text, self.page_chars)
        total_pages = len(pages)
        windows = self.create_sliding_windows(total_pages, self.window_size, self.overlap)

        chunks: list[dict] = []
        for idx, (window_start, page_indices) in enumerate(windows):
            window_text = "\n\n".join(pages[i] for i in page_indices).strip()
            if not window_text:
                continue
            page_range = f"{page_indices[0] + 1}-{page_indices[-1] + 1}"
            metadata: dict[str, Any] = {
                "window_start": window_start,
                "page_range": page_range,
                "page_number": page_indices[0] + 1,
                "total_pages": total_pages,
                "window_size": self.window_size,
                "overlap": self.overlap,
                "chunk_index": idx,
                "total_chunks": len(windows),
            }
            chunks.append({"text": window_text, "metadata": metadata})

        return chunks


def flat_chunk_text(
    text: str,
    chunk_size: int = 1024,
    overlap: int = 200,
) -> list[str]:
    """Split *text* into overlapping windows of ~*chunk_size* chars.

    Simple character-level chunker for small documents where the
    overhead of page-splitting is unnecessary.
    """
    if not text:
        return []
    if len(text) <= chunk_size:
        return [text]
    chunks: list[str] = []
    step = max(1, chunk_size - overlap)
    for start in range(0, len(text), step):
        end = start + chunk_size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk)
        if end >= len(text):
            break
    return chunks


class _ConverterResult:
    """Minimal stand-in for ``markitdown.DocumentConverterResult``.

    Avoids importing markitdown (which pulls in pandas/numpy) at runtime
    in environments where those deps are broken. The result is compatible
    with the ``text_content`` / ``markdown`` attributes that callers expect.
    """

    def __init__(self, markdown: str) -> None:
        self.markdown = markdown

    @property
    def text_content(self) -> str:
        return self.markdown

    def __str__(self) -> str:
        return self.markdown


class VisionSlidingWindowConverter(DocumentConverter):
    """PDF / DOCX / PPTX → markdown via sliding-window vision OCR.

    Renders document pages to images via PyMuPDF, groups them into sliding
    windows, and sends each combined image to a vision LLM for structured
    extraction.  Extends the MarkItDown ``DocumentConverter`` interface so
    it can be registered in the converter chain.

    For PDF files, pages are rendered directly.  For DOCX and PPTX files,
    the document is first converted to an in-memory PDF (via
    ``mammoth``+``reportlab`` for DOCX, ``python-pptx``+``reportlab``
    for PPTX), then the PDF pages are rendered as normal.

    When ``llm_client`` is not provided in *kwargs*, ``accepts()``
    returns ``False`` and MarkItDown falls back to the built-in
    converters (pdfminer for PDF, mammoth for DOCX, python-pptx for
    PPTX).
    """

    ACCEPTED_EXTENSIONS = {".pdf", ".docx", ".pptx"}
    ACCEPTED_MIME_PREFIXES = {
        "application/pdf",
        "application/x-pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }

    def __init__(
        self,
        window_size: int = SLIDING_WINDOW_SIZE,
        overlap: int = SLIDING_WINDOW_OVERLAP,
        dpi: int = 144,
    ) -> None:
        self.window_size = window_size
        self.overlap = overlap
        self.dpi = dpi

    def accepts(self, file_stream: io.BytesIO, stream_info: Any, **kwargs: Any) -> bool:
        extension = (getattr(stream_info, "extension", None) or "").lower()
        mimetype = (getattr(stream_info, "mimetype", None) or "").lower()

        if extension not in self.ACCEPTED_EXTENSIONS and not any(
            mimetype.startswith(p) for p in self.ACCEPTED_MIME_PREFIXES
        ):
            logger.debug(
                "VisionSlidingWindowConverter.accepts: declined %s (extension/mimetype not matched)",
                extension,
            )
            return False

        llm_client = kwargs.get("llm_client")
        if llm_client is not None:
            logger.debug(
                "VisionSlidingWindowConverter.accepts: accepted %s (llm_client=%s, llm_model=%s)",
                extension,
                type(llm_client).__name__,
                kwargs.get("llm_model", "<none>"),
            )
        else:
            logger.debug(
                "VisionSlidingWindowConverter.accepts: declined %s (no llm_client in kwargs)",
                extension,
            )
        return llm_client is not None

    def convert(self, file_stream: io.BytesIO, stream_info: Any, **kwargs: Any) -> Any:
        try:
            import pymupdf
        except ImportError as exc:
            raise ImportError(
                "PyMuPDF is required for VisionSlidingWindowConverter. "
                "Install it with: pip install PyMuPDF"
            ) from exc

        llm_client = kwargs.get("llm_client")
        llm_model = kwargs.get("llm_model")
        llm_prompt = kwargs.get("llm_prompt")

        extension = (getattr(stream_info, "extension", None) or "").lower()
        raw_bytes = file_stream.read()

        logger.debug(
            "VisionSlidingWindowConverter.convert: starting %s, "
            "llm_client=%s, llm_model=%s, %d bytes",
            extension,
            type(llm_client).__name__ if llm_client else "None",
            llm_model or "<none>",
            len(raw_bytes),
        )

        # Convert DOCX/PPTX to PDF first, then open with PyMuPDF.
        if extension in (".docx", ".pptx"):
            pdf_bytes = self._convert_office_to_pdf(raw_bytes, extension)
            if pdf_bytes is None:
                logger.warning(
                    "Could not convert %s to PDF for vision OCR — "
                    "raising to allow MarkItDown fallback",
                    extension,
                )
                raise RuntimeError(
                    f"Vision OCR: could not convert {extension} to PDF. "
                    f"Falling back to built-in converter."
                )
            doc = pymupdf.open(stream=pdf_bytes, filetype="pdf")
        else:
            pdf_bytes = io.BytesIO(raw_bytes)
            doc = pymupdf.open(stream=pdf_bytes, filetype="pdf")
        total_pages = len(doc)

        windows = SlidingWindowChunker.create_sliding_windows(
            total_pages, self.window_size, self.overlap
        )

        def process_window(window_start: int, page_numbers: list[int]) -> tuple[int, str]:
            try:
                combined = self._combine_page_images(doc, page_numbers)
                if not combined.getvalue():
                    return window_start, ""
                content = self._vision_ocr(
                    combined, llm_client, llm_model, llm_prompt
                )
                logger.info(
                    "Vision OCR window pages %s: %d chars",
                    [p + 1 for p in page_numbers],
                    len(content),
                )
                return window_start, content
            except Exception as exc:
                logger.error(
                    "Error processing window starting at page %d: %s",
                    window_start + 1,
                    exc,
                )
                return window_start, ""

        window_contents: dict[int, str] = {}
        with ThreadPoolExecutor(max_workers=min(MAX_WORKERS, len(windows))) as executor:
            future_to_window = {
                executor.submit(process_window, ws, pn): ws
                for ws, pn in windows
            }
            for future in as_completed(future_to_window):
                try:
                    ws, content = future.result()
                    window_contents[ws] = content
                except Exception as exc:
                    ws = future_to_window[future]
                    logger.error("Window at page %d failed: %s", ws + 1, exc)
                    window_contents[ws] = ""

        markdown_parts: list[str] = []
        for window_start, page_numbers in sorted(windows):
            content = window_contents.get(window_start, "")
            if content:
                page_range = f"{page_numbers[0] + 1}-{page_numbers[-1] + 1}"
                markdown_parts.append(f"<!-- pages {page_range} -->\n{content}")

        doc.close()
        total_chars = sum(len(p) for p in markdown_parts)
        logger.debug(
            "VisionSlidingWindowConverter.convert: completed %s, %d pages, %d windows, %d chars total, llm_model=%s",
            extension,
            total_pages,
            len(windows),
            total_chars,
            llm_model or "<none>",
        )

        if total_chars == 0:
            logger.warning(
                "VisionSlidingWindowConverter.convert: all windows produced empty "
                "results for %s (%d pages, llm_model=%s). Raising to allow "
                "MarkItDown fallback to built-in PdfConverter.",
                extension,
                total_pages,
                llm_model or "<none>",
            )
            raise RuntimeError(
                f"Vision OCR produced no content for {extension} file "
                f"({total_pages} pages, model={llm_model or 'unknown'}). "
                f"Falling back to built-in converter."
            )

        return _ConverterResult(markdown="\n\n".join(markdown_parts))

    def _combine_page_images(
        self,
        doc: Any,
        page_numbers: list[int],
    ) -> io.BytesIO:
        """Combine multiple PDF pages into a single horizontal PNG."""
        try:
            from PIL import Image
        except ImportError:
            return io.BytesIO()

        images: list[Any] = []
        total_width = 0
        max_height = 0

        for page_num in page_numbers:
            if page_num < len(doc):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(dpi=self.dpi)
                img_bytes = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_bytes))
                images.append(img)
                total_width += img.width
                max_height = max(max_height, img.height)

        if not images:
            return io.BytesIO()

        combined = Image.new("RGB", (total_width, max_height), "white")
        x_offset = 0
        for img in images:
            combined.paste(img, (x_offset, 0))
            x_offset += img.width

        stream = io.BytesIO()
        combined.save(stream, format="PNG")
        stream.seek(0)
        return stream

    @staticmethod
    def _convert_office_to_pdf(raw_bytes: bytes, extension: str) -> io.BytesIO | None:
        """Convert DOCX or PPTX bytes to an in-memory PDF.

        Uses ``mammoth`` to extract HTML from DOCX, then ``reportlab``
        to render the text into a paginated PDF.  For PPTX, uses
        ``python-pptx`` to extract slide content and ``reportlab`` to
        produce one PDF page per slide.

        Returns a ``BytesIO`` containing the PDF bytes, or ``None`` if
        the conversion fails (missing dependency, corrupt file, etc.).
        When ``None`` is returned the caller should return an empty
        result so MarkItDown falls back to its built-in converter.
        """
        if extension == ".docx":
            return VisionSlidingWindowConverter._docx_to_pdf(raw_bytes)
        if extension == ".pptx":
            return VisionSlidingWindowConverter._pptx_to_pdf(raw_bytes)
        return None

    @staticmethod
    def _docx_to_pdf(raw_bytes: bytes) -> io.BytesIO | None:
        """Convert DOCX bytes to PDF via mammoth + reportlab.

        Strategy:
        1. Use ``mammoth`` to extract clean HTML from the DOCX.
        2. Parse the HTML into plain text paragraphs.
        3. Use ``reportlab`` to render the paragraphs into a paginated
           PDF (A4 pages). This preserves the reading order and
           paragraph structure — which is exactly what the vision LLM
           needs to see.

        Embedded images in the DOCX are not rendered in this path
        (mammoth extracts them separately). For full visual fidelity,
        use LibreOffice headless (not available in all environments).
        """
        try:
            import mammoth
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib.units import mm
            from reportlab.platypus import (
                Paragraph,
                SimpleDocTemplate,
                Spacer,
            )
        except ImportError as exc:
            logger.warning(
                "mammoth or reportlab not installed — cannot convert DOCX to PDF: %s",
                exc,
            )
            return None

        try:
            # Extract HTML from DOCX
            result = mammoth.convert_to_html(io.BytesIO(raw_bytes))
            html_content = result.value

            # Strip HTML to plain text paragraphs for reportlab
            import re

            # Split on block-level tags to get paragraphs
            blocks = re.split(r"</(?:p|h[1-6]|li|div|tr|table)>", html_content)
            paragraphs: list[str] = []
            for block in blocks:
                text = re.sub(r"<[^>]+>", "", block).strip()
                if text:
                    paragraphs.append(text)

            if not paragraphs:
                return None

            # Build PDF with reportlab
            pdf_buf = io.BytesIO()
            doc = SimpleDocTemplate(
                pdf_buf,
                pagesize=A4,
                leftMargin=25 * mm,
                rightMargin=25 * mm,
                topMargin=25 * mm,
                bottomMargin=25 * mm,
            )
            styles = getSampleStyleSheet()
            story = []
            for para_text in paragraphs:
                # Escape XML special chars for reportlab Paragraph
                escaped = (
                    para_text.replace("&", "&amp;")
                    .replace("<", "&lt;")
                    .replace(">", "&gt;")
                )
                story.append(Paragraph(escaped, styles["Normal"]))
                story.append(Spacer(1, 6))

            doc.build(story)
            pdf_buf.seek(0)
            return pdf_buf

        except Exception as exc:
            logger.error("DOCX → PDF conversion failed: %s", exc)
            return None

    @staticmethod
    def _pptx_to_pdf(raw_bytes: bytes) -> io.BytesIO | None:
        """Convert PPTX bytes to PDF via python-pptx + reportlab.

        Strategy:
        1. Use ``python-pptx`` to read slide content (title, body text,
           notes).
        2. Render each slide as a PDF page via ``reportlab`` with the
           slide title as a heading and body text as paragraphs.

        This captures the textual content and structure. For full visual
        fidelity (charts, images, animations), use LibreOffice headless.
        """
        try:
            from pptx import Presentation
            from reportlab.lib.pagesizes import landscape, A4
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib.units import mm
            from reportlab.platypus import (
                PageBreak,
                Paragraph,
                SimpleDocTemplate,
                Spacer,
            )
        except ImportError as exc:
            logger.warning(
                "python-pptx or reportlab not installed — cannot convert PPTX to PDF: %s",
                exc,
            )
            return None

        try:
            prs = Presentation(io.BytesIO(raw_bytes))

            pdf_buf = io.BytesIO()
            doc = SimpleDocTemplate(
                pdf_buf,
                pagesize=landscape(A4),
                leftMargin=20 * mm,
                rightMargin=20 * mm,
                topMargin=20 * mm,
                bottomMargin=20 * mm,
            )
            styles = getSampleStyleSheet()
            story = []

            for slide_idx, slide in enumerate(prs.slides):
                if slide_idx > 0:
                    story.append(PageBreak())

                # Extract title
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        escaped = (
                            title_text.replace("&", "&amp;")
                            .replace("<", "&lt;")
                            .replace(">", "&gt;")
                        )
                        story.append(Paragraph(escaped, styles["Heading1"]))
                        story.append(Spacer(1, 12))

                # Extract body text from all shapes
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text and text != (slide.shapes.title.text.strip() if slide.shapes.title else ""):
                                escaped = (
                                    text.replace("&", "&amp;")
                                    .replace("<", "&lt;")
                                    .replace(">", "&gt;")
                                )
                                story.append(Paragraph(escaped, styles["Normal"]))
                                story.append(Spacer(1, 4))

                # Extract notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        escaped = (
                            notes_text.replace("&", "&amp;")
                            .replace("<", "&lt;")
                            .replace(">", "&gt;")
                        )
                        story.append(Spacer(1, 8))
                        story.append(
                            Paragraph(f"<i>Notes: {escaped}</i>", styles["Normal"])
                        )

            if not story:
                return None

            doc.build(story)
            pdf_buf.seek(0)
            return pdf_buf

        except Exception as exc:
            logger.error("PPTX → PDF conversion failed: %s", exc)
            return None

    @staticmethod
    def _vision_ocr(
        image_stream: io.BytesIO,
        llm_client: Any,
        llm_model: str | None,
        prompt: str | None,
    ) -> str:
        """Send a combined page image to the vision LLM for OCR.

        Supports two client types:

        1. **Project LLMClient** (``dashboard.llm_client.LLMClient``) —
           uses the async ``chat()`` method with ``ChatMessage.images``
           to pass the base64-encoded image. Works with all providers
           (OpenAI, Google/Gemini, Ollama) via ``run_sync``.

        2. **OpenAI sync SDK** (``openai.OpenAI``) — used by
           MarkItDown's built-in ``ImageConverter`` when an
           ``llm_client`` kwarg is a sync OpenAI client. Falls back to
           ``chat.completions.create`` with the data-URI format.

        When *llm_client* is None, returns "".
        """
        import base64

        if llm_client is None:
            logger.debug("_vision_ocr: skipped (no llm_client)")
            return ""

        if not prompt or not prompt.strip():
            prompt = (
                "Extract all text content from this document page image. "
                "Preserve the structure, tables, and formatting as markdown. "
                "Include any headers, footers, and captions."
            )

        cur_pos = image_stream.tell()
        try:
            b64 = base64.b64encode(image_stream.read()).decode("utf-8")
        except Exception:
            return ""
        finally:
            image_stream.seek(cur_pos)

        data_uri = f"data:image/png;base64,{b64}"

        # --- Path 1: Project LLMClient (dashboard.llm_client.LLMClient) ---
        from dashboard.llm_client import ChatMessage as _ChatMsg, LLMClient as _LLMClient

        if isinstance(llm_client, _LLMClient):
            logger.debug(
                "_vision_ocr: Path 1 (project LLMClient), model=%s, provider=%s, image=%d bytes",
                llm_model,
                getattr(llm_client, "provider", "<unknown>"),
                len(b64),
            )
            try:
                from dashboard.llm_client import run_sync as _run_sync

                messages = [
                    _ChatMsg(role="system", content="You are a document OCR assistant. Extract text from images as structured markdown."),
                    _ChatMsg(role="user", content=prompt, images=[data_uri]),
                ]
                result = _run_sync(llm_client.chat(messages))
                logger.debug(
                    "_vision_ocr: Path 1 completed, %d chars returned",
                    len(result.content or ""),
                )
                return result.content or ""
            except Exception as exc:
                logger.error("Vision LLM call (project client) failed: %s", exc)
                return ""

        # --- Path 2: OpenAI sync SDK (openai.OpenAI) ---
        if hasattr(llm_client, "chat") and hasattr(llm_client.chat, "completions"):
            model = llm_model or "gpt-4o"
            logger.debug(
                "_vision_ocr: Path 2 (OpenAI sync SDK), model=%s",
                model,
            )
            messages = [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": data_uri}},
                    ],
                }
            ]
            try:
                response = llm_client.chat.completions.create(
                    model=model, messages=messages
                )
                content = response.choices[0].message.content or ""
                logger.debug(
                    "_vision_ocr: Path 2 completed, %d chars returned",
                    len(content),
                )
                return content
            except Exception as exc:
                logger.error("Vision LLM call (OpenAI SDK) failed: %s", exc)
                return ""

        # --- Path 3: Unknown client — try generic chat with images ---
        logger.debug(
            "_vision_ocr: Path 3 (generic client), client_type=%s",
            type(llm_client).__name__,
        )
        try:
            messages = [
                _ChatMsg(role="user", content=prompt, images=[data_uri]),
            ]
            result = _run_sync(llm_client.chat(messages))
            return result.content or ""
        except Exception as exc:
            logger.error("Vision LLM call (generic) failed: %s", exc)
            return ""


__all__ = [
    "SlidingWindowChunker",
    "VisionSlidingWindowConverter",
    "flat_chunk_text",
]

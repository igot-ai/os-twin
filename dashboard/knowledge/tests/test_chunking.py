"""Unit and functional tests for dashboard.knowledge.chunking.

Covers:
- SlidingWindowChunker (unit)
- flat_chunk_text (unit)
- VisionSlidingWindowConverter (unit: accepts, _vision_ocr with mock clients)
- VisionSlidingWindowConverter (functional: _combine_page_images with real PDF)
- KnowledgeLLM.vision_ocr (unit with mocked LLMClient)
- Integration: Ingestor._parse_file with sliding-window options
"""

from __future__ import annotations

import base64
import io
import tempfile
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock, patch

import pytest

from dashboard.knowledge.chunking import (
    MAX_WORKERS,
    SlidingWindowChunker,
    VisionSlidingWindowConverter,
    flat_chunk_text,
)
from dashboard.knowledge.config import SLIDING_WINDOW_OVERLAP, SLIDING_WINDOW_SIZE


# ---------------------------------------------------------------------------
# Unit: flat_chunk_text
# ---------------------------------------------------------------------------


class TestFlatChunkText:
    def test_empty_returns_empty(self):
        assert flat_chunk_text("") == []

    def test_small_text_returns_single(self):
        assert flat_chunk_text("hello", chunk_size=1024) == ["hello"]

    def test_splits_long_text(self):
        text = "word " * 200
        chunks = flat_chunk_text(text, chunk_size=50, overlap=10)
        assert len(chunks) >= 2
        assert all(c.strip() for c in chunks)

    def test_overlap_preserved(self):
        text = "a" * 200
        chunks = flat_chunk_text(text, chunk_size=100, overlap=20)
        if len(chunks) >= 2:
            overlap_text = chunks[0][-20:]
            assert overlap_text == chunks[1][:20]

    def test_no_overlap(self):
        text = "x" * 200
        chunks = flat_chunk_text(text, chunk_size=100, overlap=0)
        assert len(chunks) == 2

    def test_whitespace_only_chunks_skipped(self):
        text = "hello" + " " * 100 + "world"
        chunks = flat_chunk_text(text, chunk_size=50, overlap=0)
        assert all(c.strip() for c in chunks)


# ---------------------------------------------------------------------------
# Unit: SlidingWindowChunker
# ---------------------------------------------------------------------------


class TestSlidingWindowChunker:
    def test_invalid_window_size(self):
        with pytest.raises(ValueError, match="window_size"):
            SlidingWindowChunker(window_size=0)

    def test_invalid_overlap_negative(self):
        with pytest.raises(ValueError, match="non-negative"):
            SlidingWindowChunker(overlap=-1)

    def test_invalid_overlap_ge_window(self):
        with pytest.raises(ValueError, match="less than window_size"):
            SlidingWindowChunker(window_size=3, overlap=3)

    def test_chunk_empty_text(self):
        chunker = SlidingWindowChunker()
        assert chunker.chunk("") == []

    def test_chunk_whitespace_only(self):
        chunker = SlidingWindowChunker()
        assert chunker.chunk("   \n\n  ") == []

    def test_chunk_small_text_single_window(self):
        chunker = SlidingWindowChunker(window_size=3, overlap=1, page_chars=1024)
        chunks = chunker.chunk("Short text")
        assert len(chunks) == 1
        assert chunks[0]["text"] == "Short text"
        assert "page_range" in chunks[0]["metadata"]
        assert chunks[0]["metadata"]["chunk_index"] == 0
        assert chunks[0]["metadata"]["total_chunks"] == 1

    def test_chunk_large_text_multiple_windows(self):
        paragraphs = [f"Paragraph {i} content." for i in range(20)]
        text = "\n\n".join(paragraphs)
        chunker = SlidingWindowChunker(window_size=3, overlap=1, page_chars=100)
        chunks = chunker.chunk(text)
        assert len(chunks) > 1
        for c in chunks:
            assert "page_range" in c["metadata"]
            assert "window_start" in c["metadata"]
            assert "total_pages" in c["metadata"]
            assert "window_size" in c["metadata"]
            assert c["metadata"]["window_size"] == 3
            assert c["metadata"]["overlap"] == 1

    def test_overlapping_windows_share_pages(self):
        chunker = SlidingWindowChunker(window_size=3, overlap=1, page_chars=50)
        # Make paragraphs long enough to each occupy its own page
        paragraphs = [f"Page {i}: " + "x" * 60 for i in range(8)]
        text = "\n\n".join(paragraphs)
        chunks = chunker.chunk(text)
        # With 8 pages, window_size=3, overlap=1:
        # Windows: [0-2], [2-4], [4-6], [6-7]
        assert chunks[0]["metadata"]["page_range"] == "1-3"
        assert chunks[1]["metadata"]["page_range"] == "3-5"

    def test_split_into_pages_paragraph_boundaries(self):
        text = "Para 1\n\nPara 2\n\nPara 3"
        pages = SlidingWindowChunker.split_into_pages(text, max_page_chars=10)
        assert len(pages) == 3
        assert "Para 1" in pages[0]
        assert "Para 2" in pages[1]

    def test_split_into_pages_long_para_hard_split(self):
        text = "x" * 500
        pages = SlidingWindowChunker.split_into_pages(text, max_page_chars=100)
        assert len(pages) >= 5

    def test_create_sliding_windows_exact_fit(self):
        windows = SlidingWindowChunker.create_sliding_windows(6, window_size=3, overlap=1)
        assert len(windows) == 3  # [0-2], [2-4], [4-5]
        assert windows[0] == (0, [0, 1, 2])
        assert windows[1] == (2, [2, 3, 4])
        assert windows[2] == (4, [4, 5])

    def test_create_sliding_windows_single_page(self):
        windows = SlidingWindowChunker.create_sliding_windows(1, window_size=3, overlap=1)
        assert len(windows) == 1
        assert windows[0] == (0, [0])


# ---------------------------------------------------------------------------
# Unit: VisionSlidingWindowConverter.accepts
# ---------------------------------------------------------------------------


class MockStreamInfo:
    def __init__(self, extension="", mimetype=""):
        self.extension = extension
        self.mimetype = mimetype


class TestVisionConverterAccepts:
    def test_accepts_pdf_with_llm_client(self):
        converter = VisionSlidingWindowConverter()
        stream = io.BytesIO(b"fake")
        info = MockStreamInfo(extension=".pdf")
        assert converter.accepts(stream, info, llm_client=MagicMock()) is True

    def test_rejects_pdf_without_llm_client(self):
        converter = VisionSlidingWindowConverter()
        stream = io.BytesIO(b"fake")
        info = MockStreamInfo(extension=".pdf")
        assert converter.accepts(stream, info) is False

    def test_accepts_docx_with_llm_client(self):
        converter = VisionSlidingWindowConverter()
        stream = io.BytesIO(b"fake")
        info = MockStreamInfo(extension=".docx")
        assert converter.accepts(stream, info, llm_client=MagicMock()) is True

    def test_accepts_pptx_with_llm_client(self):
        converter = VisionSlidingWindowConverter()
        stream = io.BytesIO(b"fake")
        info = MockStreamInfo(extension=".pptx")
        assert converter.accepts(stream, info, llm_client=MagicMock()) is True

    def test_rejects_unsupported_extension(self):
        converter = VisionSlidingWindowConverter()
        stream = io.BytesIO(b"fake")
        info = MockStreamInfo(extension=".xlsx")
        assert converter.accepts(stream, info, llm_client=MagicMock()) is False

    def test_accepts_pdf_by_mime_type(self):
        converter = VisionSlidingWindowConverter()
        stream = io.BytesIO(b"fake")
        info = MockStreamInfo(mimetype="application/pdf")
        assert converter.accepts(stream, info, llm_client=MagicMock()) is True

    def test_accepts_xpdf_by_mime_type(self):
        converter = VisionSlidingWindowConverter()
        stream = io.BytesIO(b"fake")
        info = MockStreamInfo(mimetype="application/x-pdf")
        assert converter.accepts(stream, info, llm_client=MagicMock()) is True

    def test_rejects_non_document_mime(self):
        converter = VisionSlidingWindowConverter()
        stream = io.BytesIO(b"fake")
        info = MockStreamInfo(mimetype="application/zip")
        assert converter.accepts(stream, info, llm_client=MagicMock()) is False

    def test_accepts_docx_mime_type(self):
        converter = VisionSlidingWindowConverter()
        stream = io.BytesIO(b"fake")
        info = MockStreamInfo(mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        assert converter.accepts(stream, info, llm_client=MagicMock()) is True

    def test_accepts_pptx_mime_type(self):
        converter = VisionSlidingWindowConverter()
        stream = io.BytesIO(b"fake")
        info = MockStreamInfo(mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        assert converter.accepts(stream, info, llm_client=MagicMock()) is True

    def test_rejects_no_info(self):
        converter = VisionSlidingWindowConverter()
        stream = io.BytesIO(b"fake")
        info = MockStreamInfo()
        assert converter.accepts(stream, info, llm_client=MagicMock()) is False


# ---------------------------------------------------------------------------
# Unit: VisionSlidingWindowConverter._vision_ocr with mock LLMClient
# ---------------------------------------------------------------------------


class TestVisionOCRWithProjectClient:
    """Test _vision_ocr using the project's dashboard.llm_client.LLMClient."""

    def _make_mock_llm_client(self, response_text="Extracted text from image"):
        """Build a mock LLMClient that returns a ChatMessage."""
        from dashboard.llm_client import ChatMessage, LLMClient

        mock = MagicMock(spec=LLMClient)
        mock.chat = AsyncMock(
            return_value=ChatMessage(role="assistant", content=response_text)
        )
        return mock

    def test_vision_ocr_with_project_client(self):
        mock_client = self._make_mock_llm_client("OCR result: Hello World")
        img_stream = io.BytesIO(b"\x89PNG\r\n\x1a\n" + b"\x00" * 100)
        result = VisionSlidingWindowConverter._vision_ocr(
            img_stream, mock_client, "test-model", None
        )
        assert result == "OCR result: Hello World"
        mock_client.chat.assert_called_once()

    def test_vision_ocr_with_project_client_sends_data_uri(self):
        mock_client = self._make_mock_llm_client()
        img_stream = io.BytesIO(b"\x89PNG" + b"\x00" * 50)
        VisionSlidingWindowConverter._vision_ocr(
            img_stream, mock_client, "gemini-test", "Custom prompt"
        )
        call_args = mock_client.chat.call_args
        messages = call_args[0][0]
        # Second message should have images with data URI
        user_msg = messages[1]
        assert hasattr(user_msg, "images")
        assert len(user_msg.images) == 1
        assert user_msg.images[0].startswith("data:image/png;base64,")

    def test_vision_ocr_with_project_client_none_returns_empty(self):
        img_stream = io.BytesIO(b"\x89PNG" + b"\x00" * 50)
        result = VisionSlidingWindowConverter._vision_ocr(
            img_stream, None, "test", None
        )
        assert result == ""

    def test_vision_ocr_with_project_client_exception_returns_empty(self):
        mock_client = self._make_mock_llm_client()
        mock_client.chat = AsyncMock(side_effect=RuntimeError("API error"))
        img_stream = io.BytesIO(b"\x89PNG" + b"\x00" * 50)
        result = VisionSlidingWindowConverter._vision_ocr(
            img_stream, mock_client, "test", None
        )
        assert result == ""

    def test_vision_ocr_default_prompt(self):
        mock_client = self._make_mock_llm_client()
        img_stream = io.BytesIO(b"\x89PNG" + b"\x00" * 50)
        VisionSlidingWindowConverter._vision_ocr(
            img_stream, mock_client, "test", None
        )
        call_args = mock_client.chat.call_args
        messages = call_args[0][0]
        user_msg = messages[1]
        assert "Extract all text content" in user_msg.content

    def test_vision_ocr_custom_prompt(self):
        mock_client = self._make_mock_llm_client()
        img_stream = io.BytesIO(b"\x89PNG" + b"\x00" * 50)
        VisionSlidingWindowConverter._vision_ocr(
            img_stream, mock_client, "test", "Describe this chart"
        )
        call_args = mock_client.chat.call_args
        messages = call_args[0][0]
        user_msg = messages[1]
        assert user_msg.content == "Describe this chart"


# ---------------------------------------------------------------------------
# Unit: VisionSlidingWindowConverter._vision_ocr with mock OpenAI sync client
# ---------------------------------------------------------------------------


class TestVisionOCRWithOpenAISyncClient:
    """Test _vision_ocr using a mock openai.OpenAI sync client."""

    def _make_mock_openai_client(self, response_text="OpenAI OCR result"):
        mock_client = MagicMock()
        mock_choice = MagicMock()
        mock_choice.message.content = response_text
        mock_response = MagicMock()
        mock_response.choices = [mock_choice]
        mock_client.chat.completions.create = MagicMock(return_value=mock_response)
        return mock_client

    def test_vision_ocr_with_openai_sync_client(self):
        mock_client = self._make_mock_openai_client("Hello from OpenAI")
        img_stream = io.BytesIO(b"\x89PNG" + b"\x00" * 50)
        result = VisionSlidingWindowConverter._vision_ocr(
            img_stream, mock_client, "gpt-4o", None
        )
        assert result == "Hello from OpenAI"

    def test_vision_ocr_openai_sends_data_uri_in_messages(self):
        mock_client = self._make_mock_openai_client()
        img_stream = io.BytesIO(b"\x89PNG" + b"\x00" * 50)
        VisionSlidingWindowConverter._vision_ocr(
            img_stream, mock_client, "gpt-4o", "Read this"
        )
        call_kwargs = mock_client.chat.completions.create.call_args
        messages = call_kwargs[1]["messages"]
        assert messages[0]["role"] == "user"
        content = messages[0]["content"]
        assert content[0]["type"] == "text"
        assert content[0]["text"] == "Read this"
        assert content[1]["type"] == "image_url"
        assert content[1]["image_url"]["url"].startswith("data:image/png;base64,")


# ---------------------------------------------------------------------------
# Functional: VisionSlidingWindowConverter._combine_page_images
# ---------------------------------------------------------------------------


class TestCombinePageImages:
    """Test _combine_page_images with a real PDF (created via PyMuPDF)."""

    @pytest.fixture
    def sample_pdf_doc(self):
        try:
            import pymupdf
        except ImportError:
            pytest.skip("PyMuPDF not installed")
        doc = pymupdf.open()
        for i in range(5):
            page = doc.new_page(width=612, height=792)
            page.insert_text((72, 72), f"Page {i + 1} content", fontsize=14)
        return doc

    def test_combine_single_page(self, sample_pdf_doc):
        converter = VisionSlidingWindowConverter(dpi=72)
        result = converter._combine_page_images(sample_pdf_doc, [0])
        assert result.getvalue()
        # Should be a valid PNG
        assert result.getvalue()[:4] == b"\x89PNG"

    def test_combine_multiple_pages(self, sample_pdf_doc):
        converter = VisionSlidingWindowConverter(dpi=72)
        result = converter._combine_page_images(sample_pdf_doc, [0, 1, 2])
        assert result.getvalue()
        from PIL import Image
        img = Image.open(result)
        assert img.width > 0
        assert img.height > 0

    def test_combine_empty_page_list(self, sample_pdf_doc):
        converter = VisionSlidingWindowConverter(dpi=72)
        result = converter._combine_page_images(sample_pdf_doc, [])
        assert result.getvalue() == b""

    def test_combine_out_of_range_pages(self, sample_pdf_doc):
        converter = VisionSlidingWindowConverter(dpi=72)
        result = converter._combine_page_images(sample_pdf_doc, [0, 99])
        assert result.getvalue()


# ---------------------------------------------------------------------------
# Functional: VisionSlidingWindowConverter.convert end-to-end with mock LLM
# ---------------------------------------------------------------------------


class TestVisionConverterConvert:
    """Test the full convert() flow with a real PDF and mocked LLM."""

    @pytest.fixture
    def sample_pdf_bytes(self):
        try:
            import pymupdf
        except ImportError:
            pytest.skip("PyMuPDF not installed")
        doc = pymupdf.open()
        for i in range(5):
            page = doc.new_page(width=612, height=792)
            page.insert_text((72, 72), f"This is page {i + 1} of the test document.", fontsize=14)
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        doc.close()
        return buf

    def test_convert_produces_markdown(self, sample_pdf_bytes):
        from dashboard.llm_client import ChatMessage, LLMClient

        mock_client = MagicMock(spec=LLMClient)
        mock_client.chat = AsyncMock(
            return_value=ChatMessage(role="assistant", content="Extracted page content")
        )

        converter = VisionSlidingWindowConverter(window_size=3, overlap=1, dpi=72)
        stream_info = MockStreamInfo(extension=".pdf")
        result = converter.convert(
            sample_pdf_bytes, stream_info,
            llm_client=mock_client, llm_model="test-model"
        )
        assert result.markdown
        assert "Extracted page content" in result.markdown
        # Should have page-range comments
        assert "<!-- pages" in result.markdown

    def test_convert_empty_pdf(self):
        try:
            import pymupdf
        except ImportError:
            pytest.skip("PyMuPDF not installed")
        # Create a PDF with 1 blank page (PyMuPDF can't save zero-page PDFs)
        doc = pymupdf.open()
        doc.new_page(width=612, height=792)
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        doc.close()

        from dashboard.llm_client import ChatMessage, LLMClient

        mock_client = MagicMock(spec=LLMClient)
        mock_client.chat = AsyncMock(
            return_value=ChatMessage(role="assistant", content="")
        )

        converter = VisionSlidingWindowConverter(window_size=3, overlap=1, dpi=72)
        stream_info = MockStreamInfo(extension=".pdf")
        result = converter.convert(
            buf, stream_info,
            llm_client=mock_client, llm_model="test-model"
        )
        # Single blank page → empty or minimal content
        assert isinstance(result.markdown, str)


# ---------------------------------------------------------------------------
# Unit: KnowledgeLLM.vision_ocr
# ---------------------------------------------------------------------------


class TestKnowledgeLLMVisionOCR:
    """Test the vision_ocr method on KnowledgeLLM with mocked client."""

    def _make_llm(self, model="gemini-3.1-flash-lite-preview", provider="google"):
        from dashboard.knowledge.llm import KnowledgeLLM
        llm = KnowledgeLLM(api_key="fake-key", model=model, provider=provider)
        return llm

    def test_vision_ocr_unavailable_returns_empty(self):
        from dashboard.knowledge.llm import KnowledgeLLM
        llm = KnowledgeLLM(model="", provider=None)
        assert llm.vision_ocr("data:image/png;base64,abc") == ""

    def test_vision_ocr_sends_image_via_chat_message(self):
        from dashboard.llm_client import ChatMessage

        llm = self._make_llm()
        mock_client = MagicMock()
        mock_client.chat = AsyncMock(
            return_value=ChatMessage(role="assistant", content="OCR text here")
        )
        with patch.object(llm, "_get_client", return_value=mock_client):
            result = llm.vision_ocr("data:image/png;base64,abc123")
        assert result == "OCR text here"
        mock_client.chat.assert_called_once()
        messages = mock_client.chat.call_args[0][0]
        # System + user message
        assert len(messages) == 2
        user_msg = messages[1]
        assert user_msg.images == ["data:image/png;base64,abc123"]

    def test_vision_ocr_custom_prompt(self):
        from dashboard.llm_client import ChatMessage

        llm = self._make_llm()
        mock_client = MagicMock()
        mock_client.chat = AsyncMock(
            return_value=ChatMessage(role="assistant", content="Chart data")
        )
        with patch.object(llm, "_get_client", return_value=mock_client):
            result = llm.vision_ocr(
                "data:image/png;base64,abc", prompt="Describe the chart"
            )
        assert result == "Chart data"
        messages = mock_client.chat.call_args[0][0]
        assert messages[1].content == "Describe the chart"

    def test_vision_ocr_exception_returns_empty(self):
        llm = self._make_llm()
        mock_client = MagicMock()
        mock_client.chat = AsyncMock(side_effect=RuntimeError("API down"))
        with patch.object(llm, "_get_client", return_value=mock_client):
            result = llm.vision_ocr("data:image/png;base64,abc")
        assert result == ""

    def test_vision_ocr_with_vertex_model(self):
        from dashboard.llm_client import ChatMessage

        llm = self._make_llm(
            model="google-vertext/gemini-3.1-flash-lite-preview",
            provider="google-vertex",
        )
        mock_client = MagicMock()
        mock_client.chat = AsyncMock(
            return_value=ChatMessage(role="assistant", content="Vertex OCR result")
        )
        with patch.object(llm, "_get_client", return_value=mock_client):
            result = llm.vision_ocr("data:image/png;base64,abc")
        assert result == "Vertex OCR result"

    def test_create_vision_client_unavailable(self):
        from dashboard.knowledge.llm import KnowledgeLLM
        llm = KnowledgeLLM(model="", provider=None)
        assert llm.create_vision_client() is None

    def test_create_vision_client_returns_client(self):
        from dashboard.llm_client import LLMClient

        llm = self._make_llm()
        mock_client = MagicMock(spec=LLMClient)
        with patch.object(llm, "_get_client", return_value=mock_client):
            result = llm.create_vision_client()
        assert result is mock_client


# ---------------------------------------------------------------------------
# Integration: IngestOptions with sliding window
# ---------------------------------------------------------------------------


class TestIngestOptionsSlidingWindow:
    def test_default_values(self):
        from dashboard.knowledge.ingestion import IngestOptions

        opts = IngestOptions()
        assert opts.vision_ocr is True
        assert opts.sliding_window_size == 3
        assert opts.sliding_window_overlap == 1
        assert opts.vision_ocr_dpi == 144

    def test_custom_values(self):
        from dashboard.knowledge.ingestion import IngestOptions

        opts = IngestOptions(
            vision_ocr=False,
            sliding_window_size=5,
            sliding_window_overlap=2,
            vision_ocr_dpi=200,
        )
        assert opts.vision_ocr is False
        assert opts.sliding_window_size == 5
        assert opts.sliding_window_overlap == 2
        assert opts.vision_ocr_dpi == 200

    def test_sliding_window_inherits_from_config(self):
        from dashboard.knowledge.ingestion import IngestOptions

        opts = IngestOptions()
        assert opts.sliding_window_size == SLIDING_WINDOW_SIZE
        assert opts.sliding_window_overlap == SLIDING_WINDOW_OVERLAP


# ---------------------------------------------------------------------------
# Integration: Ingestor._parse_file with sliding window
# ---------------------------------------------------------------------------


class TestIngestorSlidingWindowIntegration:
    @pytest.fixture
    def make_ingestor(self):
        from dashboard.knowledge.ingestion import Ingestor, FileEntry, IngestOptions
        from dashboard.knowledge.namespace import NamespaceManager

        nm = NamespaceManager(base_dir=Path(tempfile.mkdtemp()))

        embedder = MagicMock()
        embedder.dimension.return_value = 1024

        def _graph_index_factory(ns):
            mock_idx = MagicMock()

            def _insert_nodes(nodes, **kwargs):
                from llama_index.core.graph_stores.types import KG_NODES_KEY, KG_RELATIONS_KEY
                for n in nodes:
                    n.metadata[KG_NODES_KEY] = []
                    n.metadata[KG_RELATIONS_KEY] = []

            mock_idx.insert_nodes = _insert_nodes
            return mock_idx

        class FakeStore:
            def __init__(self):
                self.chunks = []

            def add_chunks(self, chunks):
                self.chunks.extend(chunks)
                return len(chunks)

            def has_file_hash(self, h):
                return any(c["metadata"].get("file_hash") == h for c in self.chunks)

            def count_by_file_hash(self, h):
                return sum(1 for c in self.chunks if c["metadata"].get("file_hash") == h)

            def delete_by_file_hash(self, h):
                before = len(self.chunks)
                self.chunks = [c for c in self.chunks if c["metadata"].get("file_hash") != h]
                return before - len(self.chunks)

        store = FakeStore()

        ing = Ingestor(
            namespace_manager=nm,
            embedder=embedder,
            graph_index_factory=_graph_index_factory,
        )
        ing._get_store = lambda ns: store
        return ing, store

    def test_parse_small_txt_flat_chunking(self, make_ingestor):
        from dashboard.knowledge.ingestion import IngestOptions, FileEntry

        ing, store = make_ingestor
        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False) as f:
            f.write("Hello world")
            f.flush()
            fe = FileEntry(
                path=f.name, size=11, mtime=0.0, extension=".txt", content_hash="h1"
            )
        chunks = ing._parse_file(fe, IngestOptions())
        assert len(chunks) >= 1
        assert "page_range" not in chunks[0]["metadata"]

    def test_parse_large_txt_sliding_window(self, make_ingestor):
        from dashboard.knowledge.ingestion import IngestOptions, FileEntry

        ing, store = make_ingestor
        paragraphs = [f"Section {i}: " + "content " * 100 for i in range(20)]
        large_text = "\n\n".join(paragraphs)
        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False) as f:
            f.write(large_text)
            f.flush()
            fe = FileEntry(
                path=f.name,
                size=len(large_text),
                mtime=0.0,
                extension=".txt",
                content_hash="h2",
            )
        opts = IngestOptions(
            chunk_size=1024,
            sliding_window_size=3,
            sliding_window_overlap=1,
        )
        chunks = ing._parse_file(fe, opts)
        assert len(chunks) >= 1
        # Large file should use sliding window — page_range metadata present
        sw_chunks = [c for c in chunks if "page_range" in c["metadata"]]
        assert len(sw_chunks) > 0, "Expected sliding-window metadata for large file"

    def test_parse_with_vision_ocr_disabled(self, make_ingestor):
        from dashboard.knowledge.ingestion import IngestOptions, FileEntry

        ing, store = make_ingestor
        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False) as f:
            f.write("Simple content")
            f.flush()
            fe = FileEntry(
                path=f.name, size=14, mtime=0.0, extension=".txt", content_hash="h3"
            )
        opts = IngestOptions(vision_ocr=False)
        chunks = ing._parse_file(fe, opts)
        assert len(chunks) >= 1


# ---------------------------------------------------------------------------
# Unit: VisionSlidingWindowConverter._convert_office_to_pdf
# ---------------------------------------------------------------------------


def _create_minimal_docx() -> bytes:
    """Create a minimal valid DOCX file in memory."""
    import zipfile

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>''')
        zf.writestr('_rels/.rels', '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>''')
        zf.writestr('word/_rels/document.xml.rels', '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
</Relationships>''')
        zf.writestr('word/document.xml', '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    <w:p><w:r><w:t>First paragraph of the document.</w:t></w:r></w:p>
    <w:p><w:r><w:t>Second paragraph with more content.</w:t></w:r></w:p>
    <w:p><w:r><w:t>Third paragraph about important topics.</w:t></w:r></w:p>
  </w:body>
</w:document>''')
    return buf.getvalue()


def _create_minimal_pptx() -> bytes:
    """Create a minimal valid PPTX file in memory."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + content
    slide1.shapes.title.text = "Welcome to the Presentation"
    slide1.placeholders[1].text = "This is the first slide content."

    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Key Findings"
    slide2.placeholders[1].text = "Our research shows significant results."

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestDocxToPdfConversion:
    """Test _convert_office_to_pdf for DOCX files."""

    def test_docx_to_pdf_produces_valid_pdf(self):
        docx_bytes = _create_minimal_docx()
        result = VisionSlidingWindowConverter._convert_office_to_pdf(docx_bytes, ".docx")
        assert result is not None
        # Should be a valid PDF (starts with %PDF)
        pdf_data = result.read(5)
        assert pdf_data == b"%PDF-"

    def test_docx_to_pdf_has_multiple_pages(self):
        docx_bytes = _create_minimal_docx()
        result = VisionSlidingWindowConverter._convert_office_to_pdf(docx_bytes, ".docx")
        assert result is not None
        result.seek(0)
        try:
            import pymupdf
            doc = pymupdf.open(stream=result, filetype="pdf")
            assert len(doc) >= 1
            # The text should contain our paragraphs
            full_text = "".join(page.get_text() for page in doc)
            assert "First paragraph" in full_text or "paragraph" in full_text.lower()
            doc.close()
        except ImportError:
            pytest.skip("PyMuPDF not installed")

    def test_docx_to_pdf_with_invalid_bytes_returns_none(self):
        result = VisionSlidingWindowConverter._convert_office_to_pdf(b"not a docx", ".docx")
        assert result is None

    def test_docx_to_pdf_empty_docx(self):
        # Create a minimal DOCX with no paragraphs
        import zipfile
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
            zf.writestr("[Content_Types].xml", '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>''')
            zf.writestr('_rels/.rels', '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>''')
            zf.writestr('word/document.xml', '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:body></w:body></w:document>''')
        result = VisionSlidingWindowConverter._convert_office_to_pdf(buf.getvalue(), ".docx")
        # Empty DOCX should return None (no paragraphs to render)
        assert result is None

    def test_convert_office_unknown_extension_returns_none(self):
        result = VisionSlidingWindowConverter._convert_office_to_pdf(b"fake", ".xlsx")
        assert result is None


class TestPptxToPdfConversion:
    """Test _convert_office_to_pdf for PPTX files."""

    def test_pptx_to_pdf_produces_valid_pdf(self):
        pptx_bytes = _create_minimal_pptx()
        result = VisionSlidingWindowConverter._convert_office_to_pdf(pptx_bytes, ".pptx")
        assert result is not None
        pdf_data = result.read(5)
        assert pdf_data == b"%PDF-"

    def test_pptx_to_pdf_has_slides_as_pages(self):
        pptx_bytes = _create_minimal_pptx()
        result = VisionSlidingWindowConverter._convert_office_to_pdf(pptx_bytes, ".pptx")
        assert result is not None
        result.seek(0)
        try:
            import pymupdf
            doc = pymupdf.open(stream=result, filetype="pdf")
            # Should have 2 pages (one per slide)
            assert len(doc) == 2
            # Check that slide titles appear in the PDF
            full_text = "".join(page.get_text() for page in doc)
            assert "Welcome" in full_text or "Presentation" in full_text
            doc.close()
        except ImportError:
            pytest.skip("PyMuPDF not installed")

    def test_pptx_to_pdf_with_invalid_bytes_returns_none(self):
        result = VisionSlidingWindowConverter._convert_office_to_pdf(b"not a pptx", ".pptx")
        assert result is None


# ---------------------------------------------------------------------------
# Functional: VisionSlidingWindowConverter.convert with DOCX/PPTX
# ---------------------------------------------------------------------------


class TestVisionConverterDocxPptx:
    """Test the full convert() flow with DOCX/PPTX files and mocked LLM."""

    def test_convert_docx_produces_markdown(self):
        try:
            import pymupdf
        except ImportError:
            pytest.skip("PyMuPDF not installed")

        from dashboard.llm_client import ChatMessage, LLMClient

        mock_client = MagicMock(spec=LLMClient)
        mock_client.chat = AsyncMock(
            return_value=ChatMessage(role="assistant", content="Extracted DOCX content")
        )

        docx_bytes = _create_minimal_docx()
        converter = VisionSlidingWindowConverter(window_size=3, overlap=1, dpi=72)
        stream_info = MockStreamInfo(extension=".docx")
        result = converter.convert(
            io.BytesIO(docx_bytes), stream_info,
            llm_client=mock_client, llm_model="test-model"
        )
        assert result.markdown
        assert "Extracted DOCX content" in result.markdown

    def test_convert_pptx_produces_markdown(self):
        try:
            import pymupdf
        except ImportError:
            pytest.skip("PyMuPDF not installed")

        from dashboard.llm_client import ChatMessage, LLMClient

        mock_client = MagicMock(spec=LLMClient)
        mock_client.chat = AsyncMock(
            return_value=ChatMessage(role="assistant", content="Extracted PPTX content")
        )

        pptx_bytes = _create_minimal_pptx()
        converter = VisionSlidingWindowConverter(window_size=3, overlap=1, dpi=72)
        stream_info = MockStreamInfo(extension=".pptx")
        result = converter.convert(
            io.BytesIO(pptx_bytes), stream_info,
            llm_client=mock_client, llm_model="test-model"
        )
        assert result.markdown
        assert "Extracted PPTX content" in result.markdown

    def test_convert_docx_conversion_failure_returns_empty(self):
        """When DOCX-to-PDF conversion fails, convert returns empty markdown."""
        converter = VisionSlidingWindowConverter(window_size=3, overlap=1, dpi=72)
        stream_info = MockStreamInfo(extension=".docx")
        # Pass invalid DOCX bytes
        result = converter.convert(
            io.BytesIO(b"not a docx"), stream_info,
            llm_client=MagicMock(), llm_model="test-model"
        )
        # Should gracefully return empty markdown
        assert result.markdown == ""

    def test_convert_pptx_conversion_failure_returns_empty(self):
        """When PPTX-to-PDF conversion fails, convert returns empty markdown."""
        converter = VisionSlidingWindowConverter(window_size=3, overlap=1, dpi=72)
        stream_info = MockStreamInfo(extension=".pptx")
        # Pass invalid PPTX bytes
        result = converter.convert(
            io.BytesIO(b"not a pptx"), stream_info,
            llm_client=MagicMock(), llm_model="test-model"
        )
        assert result.markdown == ""
